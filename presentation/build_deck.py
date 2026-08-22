"""Build the 45-minute comfort-zone-boundary talk on the Chalmers template.

    python presentation/build_deck.py [--out presentation/czb_talk.pptx]

Audience: mixed — traffic-safety analytics people (coding/ML, some human
factors) and human-factors people (little in-depth math or coding). The deck
alternates depth deliberately: conceptual slides carry everyone, and the
math/implementation slides are flagged in the speaker notes so the HF half can
be told what to take from them without following the detail.

Every slide carries speaker notes with what to say and a rough time budget.
Content comes from notes/02, notes/04, notes/05 and README.md; figures from
figures/. Regenerate after any of those change.
"""

from __future__ import annotations

import argparse
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
REPO = HERE.parent
FIGS = REPO / "figures"
TEMPLATE_POTX = HERE / "chalmers-tekniska-ho-gskola-sv.potx"
TEMPLATE_PPTX = HERE / "_chalmers-template.pptx"

PURPLE = RGBColor(0x47, 0x2C, 0xBE)   # accent1
LILAC = RGBColor(0x67, 0x46, 0xEB)    # accent3
BLUE = RGBColor(0x36, 0xB7, 0xF6)     # accent5
TEAL = RGBColor(0x61, 0xE9, 0xD2)     # accent6
PINK = RGBColor(0xD9, 0x87, 0xBA)     # accent4
INK = RGBColor(0x22, 0x22, 0x22)      # dk1
BEIGE = RGBColor(0xF0, 0xED, 0xE6)    # lt2
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x5A, 0x5A, 0x5A)

FONT = "Arial"
L_TITLE = 0     # Titelsida
L_SECTION = 2   # Kapitel
L_HEAD = 30     # Endast rubrik
L_CLOSE = 32    # Slutsida


def ensure_template() -> Path:
    if TEMPLATE_PPTX.exists() and TEMPLATE_PPTX.stat().st_mtime >= TEMPLATE_POTX.stat().st_mtime:
        return TEMPLATE_PPTX
    with zipfile.ZipFile(TEMPLATE_POTX) as zin, \
            zipfile.ZipFile(TEMPLATE_PPTX, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(b"presentationml.template.main+xml",
                                    b"presentationml.presentation.main+xml")
            zout.writestr(item, data)
    return TEMPLATE_PPTX


def drop_existing_slides(prs) -> None:
    id_list = prs.slides._sldIdLst
    for slide_id in list(id_list):
        prs.part.drop_rel(slide_id.rId)
        id_list.remove(slide_id)


def clean(slide) -> None:
    """Remove every unfilled placeholder so nothing ships as a prompt."""
    for shape in list(slide.placeholders):
        if shape.has_text_frame:
            if not shape.text_frame.text.strip():
                shape._element.getparent().remove(shape._element)
        else:
            try:
                _ = shape.image
            except Exception:
                shape._element.getparent().remove(shape._element)


def _plain(shape, colour) -> None:
    if colour is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False


def text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.06):
    """Each run: (text, size, colour, bold[, space_before_pt])."""
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    for index, run in enumerate(runs):
        body, size, colour, bold = run[:4]
        before = run[4] if len(run) > 4 else 0
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        if before:
            para.space_before = Pt(before)
        piece = para.add_run()
        piece.text = body
        piece.font.size = Pt(size)
        piece.font.color.rgb = colour
        piece.font.bold = bold
        piece.font.name = FONT
    return box


def bullets(items, size=13, gap=7):
    """items: str, or (label, body) rendered as coloured lead-in + body."""
    runs = []
    for i, item in enumerate(items):
        before = 0 if i == 0 else gap
        if isinstance(item, tuple):
            runs.append((item[0], size, PURPLE, True, before))
            runs.append((item[1], size, INK, False, 1))
        else:
            runs.append(("–  " + item, size, INK, False, before))
    return runs


def column(slide, x, y, w, h, heading, items, *, rule=PURPLE, size=11.5, head_size=13):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.62), Pt(3.5))
    _plain(bar, rule)
    text(slide, x, y + Inches(0.12), w, Inches(0.34), [(heading, head_size, INK, True)])
    runs = []
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            runs.append((item[0], size, rule if rule != TEAL else PURPLE, True, 0 if i == 0 else 7))
            runs.append((item[1], size, INK, False, 1))
        else:
            runs.append(("–  " + item, size, INK, False, 0 if i == 0 else 5))
    return text(slide, x, y + Inches(0.52), w, h, runs, spacing=1.1)


def panel(slide, x, y, w, h, colour=BEIGE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.045
    _plain(shp, colour)
    return shp


def picture(slide, name, x, y, w=None, h=None):
    """name is resolved against figures/ first, then the repo root."""
    path = FIGS / name
    if not path.exists():
        path = REPO / name
    if not path.exists():
        print("MISSING FIGURE:", name)
        return None
    pic = slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    pic.line.color.rgb = RGBColor(0xD5, 0xD1, 0xC8)
    pic.line.width = Pt(0.75)
    return pic


def table(slide, x, y, col_w, rows, *, size=11, row_h=0.34):
    """Light hand-drawn table: rows[0] is the header. col_w in inches."""
    yy = y
    for r, row in enumerate(rows):
        head = r == 0
        if head:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x, yy + Inches(row_h) - Pt(3),
                Inches(sum(col_w)), Pt(1.6))
            _plain(bar, PURPLE)
        xx = x
        for c, cell in enumerate(row):
            text(slide, xx, yy, Inches(col_w[c]) - Inches(0.12), Inches(row_h),
                 [(str(cell), size, PURPLE if head else INK, head or c == 0)])
            xx += Inches(col_w[c])
        yy += Inches(row_h + (0.04 if head else 0.0))
    return yy


def notes(slide, body: str) -> None:
    slide.notes_slide.notes_text_frame.text = body.strip()


def head_slide(prs, title_text, *, kicker=None):
    slide = prs.slides.add_slide(prs.slide_layouts[L_HEAD])
    title = slide.shapes.title
    title.text_frame.text = ""
    para = title.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = title_text
    run.font.size = Pt(23)
    run.font.bold = True
    run.font.color.rgb = INK
    run.font.name = FONT
    if kicker:
        text(slide, Inches(0.55), Inches(0.18), Inches(9), Inches(0.3),
             [(kicker.upper(), 10.5, PURPLE, True)])
    clean(slide)
    return slide


def blank_slide(prs):
    """Layout 30 with the title stripped — a clean canvas we own entirely."""
    slide = prs.slides.add_slide(prs.slide_layouts[L_HEAD])
    for shape in list(slide.placeholders):
        shape._element.getparent().remove(shape._element)
    return slide


def section_slide(prs, number, title_text, note_text):
    """Hand-drawn divider: the Kapitel layout renders empty, so paint our own."""
    slide = blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    _plain(bg, PURPLE)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.55),
                                 Inches(0.9), Pt(4.5))
    _plain(bar, TEAL)
    text(slide, Inches(0.9), Inches(2.0), Inches(6), Inches(0.4),
         [("PART " + str(number), 15, TEAL, True)])
    text(slide, Inches(0.9), Inches(2.95), Inches(11.4), Inches(2.4),
         [(title_text, 33, WHITE, True)], spacing=1.1)
    notes(slide, note_text)
    return slide


def build(out: Path) -> None:
    prs = Presentation(str(ensure_template()))
    drop_existing_slides(prs)

    # ---- 1 - title (hand-drawn: the Titelsida layout renders empty) ----------
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    _plain(bg, PURPLE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.35),
                             Inches(0.9), Pt(4.5))
    _plain(bar, TEAL)
    text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.0),
         [("Driver comfort-zone boundaries from active inference", 34, WHITE, True)],
         spacing=1.1)
    text(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(1.4), [
        ("One scalar field instead of one indicator per scenario", 17, TEAL, False),
        ("Jonas Bärgman  ·  Chalmers University of Technology", 14, WHITE, False, 10),
    ], spacing=1.1)
    notes(s, """
Total 45 min: problem 6, active inference 12, replication 7, the method 15, validation plan 4, wrap-up 1.
If running long, the surprise-library slide in part 3 is the one to compress.
Say: this talk is about an old construct - the driver's comfort zone - and a new way to compute it.
The work builds on a 2026 Nature Communications model of collision avoidance from Waymo and TU Delft,
but replicating that model is the means, not the end. Two audiences today: I will flag the slides where
the math gets dense, and I promise there is a human-factors payoff on the other side of each of them.
""")

    # ---- 2 - the claim --------------------------------------------------------
    s = head_slide(prs, "The claim in one sentence")
    panel(s, Inches(0.55), Inches(1.45), Inches(12.2), Inches(1.55))
    text(s, Inches(0.95), Inches(1.68), Inches(11.4), Inches(1.1),
         [("A driver's comfort-zone boundary is a level set of one scalar field, computable from "
           "a preference function and recorded kinematics alone — the same in every scenario.",
           17, PURPLE, True)], spacing=1.15)
    text(s, Inches(0.55), Inches(3.45), Inches(12.2), Inches(3.2), bullets([
        ("Today: ", "comfort-zone boundaries are quantified per scenario, per kinematic indicator "
         "(min TTC here, lateral clearance there, THW elsewhere)."),
        ("The proposal: ", "replace the per-scenario indicators with one field derived from an "
         "active-inference driver model; classic indicators become projections of it."),
        ("Status: ", "the method is built, cross-checked, and has passed its first end-to-end test "
         "— it recovers the reference model's response onsets from kinematics alone. No human "
         "data has touched it yet. The human validation study is the last part of this talk."),
    ], size=14, gap=10))
    notes(s, """
~2 min. State the claim, then immediately state the epistemic status: this is a construction, not yet a
finding. Being upfront about that sets the tone and buys credibility with both audiences. The analytics
people will want the field definition (coming in part 4); the HF people will recognize the comfort-zone
construct (coming right now, part 1).
""")

    # ---- PART 1 ---------------------------------------------------------------
    section_slide(prs, 1, "Comfort zones, and the problem with how we measure them", """
~30 s. Part 1 is home turf for the human-factors half - the analytics half gets the construct they may
not know they have been modeling around.
""")

    # ---- 3 - the construct -----------------------------------------------------
    s = head_slide(prs, "The comfort zone is a fifty-year-old construct",
                   kicker="Part 1 · The problem")
    text(s, Inches(0.55), Inches(1.5), Inches(6.3), Inches(5.3), bullets([
        ("Näätänen & Summala (1976): ", "drivers regulate to keep subjective risk at "
         "zero — they act when a situation threatens to leave the zone where they feel "
         "comfortable."),
        ("Comfort-zone boundary: ", "the limit drivers do not cross voluntarily without extra "
         "motives."),
        ("Dread-zone boundary: ", "the further limit they do not cross even with extra motives."),
        ("Extra motives: ", "being late, angry, pressured — what makes a driver accept normally "
         "unacceptable discomfort."),
        ("Chalmers work: ", "LTAP/OD comfort- and dread-zone quantification (2015); pedestrian "
         "overtaking (2019); Engström et al. (2018) Great expectations gives the "
         "predictive-processing framing this talk continues."),
    ], size=13, gap=10))
    column(s, Inches(7.3), Inches(1.5), Inches(5.4), Inches(4.5),
           "How boundaries are measured today", [
               "Pick the kinematic indicator that seems to matter in that scenario",
               "Collect accepted values from field or naturalistic data",
               "Report a quantile of what drivers accept",
               ("Result: ", "a THW boundary for following, a TTC boundary for LTAP/OD, a clearance "
                "boundary for overtaking — each valid, none comparable"),
           ], rule=BLUE, size=12)
    notes(s, """
~3 min. HF half: nothing new, nod along. Analytics half: this is the construct - drivers maintain a
dynamic safety envelope and act when it is threatened, and the boundary moves with motivation. Mention
your own stake: co-author of Great expectations, and the LTAP/OD and overtaking quantification is
Chalmers work - so this project is a continuation, not a pivot.
""")

    # ---- 4 - three costs --------------------------------------------------------
    s = head_slide(prs, "Per-scenario indicators carry three costs",
                   kicker="Part 1 · The problem")
    ys = Inches(1.55)
    for i, (t1, t2, colr) in enumerate([
        ("Boundaries are not comparable across scenarios",
         "A THW boundary and a lateral-clearance boundary are different objects. There is no way to "
         "say whether a driver is equally close to their boundary in two scenarios.", BLUE),
        ("The choice of indicator is a modeling assumption in disguise",
         "Nothing tells you whether min TTC, required deceleration or PET is the variable the driver "
         "is actually regulating.", PINK),
        ("No mechanism",
         "A quantile of accepted TTC describes behavior. It does not explain why drivers act when "
         "they do, and it cannot predict the boundary in a scenario not yet measured.", PURPLE),
    ]):
        panel(s, Inches(0.55), ys, Inches(12.2), Inches(1.55))
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), ys, Inches(0.09), Inches(1.55))
        _plain(bar, colr)
        text(s, Inches(0.95), ys + Inches(0.18), Inches(11.5), Inches(1.2),
             [(str(i + 1) + ".  " + t1, 15, INK, True),
              (t2, 12.5, GREY, False, 4)], spacing=1.1)
        ys += Inches(1.78)
    notes(s, """
~2.5 min. The third cost is the one to lean on: description versus mechanism. A quantile summarizes; it
cannot transfer. If someone asks 'but the indicators work fine' - agree: they do, per scenario. The point
is what they cannot do: compare across scenarios and predict unmeasured ones. That is the gap the rest
of the talk fills.
""")

    # ---- PART 2 -----------------------------------------------------------------
    section_slide(prs, 2, "Active inference, as instantiated for driving", """
~30 s. Signpost honestly: this part has the steepest math of the talk. Two conceptual slides everyone
should follow, then two implementation slides aimed at the analytics half - tell the HF half what to
take from each as you go.
""")

    # ---- 5 - core idea ------------------------------------------------------------
    s = head_slide(prs, "One idea: perception and action both minimize surprise",
                   kicker="Part 2 · Active inference")
    text(s, Inches(0.55), Inches(1.5), Inches(6.3), Inches(5.4), bullets([
        ("From computational neuroscience ", "(Friston's free-energy principle): the brain holds a "
         "generative model of the world and works to reduce the mismatch between what it predicts "
         "and what it senses."),
        ("Perception ", "updates beliefs so predictions fit the senses (state estimation)."),
        ("Action ", "changes the world so the senses fit the predictions (control)."),
        ("The trick is the preference prior p(o): ", "the observations the driver wants are the "
         "ones the model treats as most probable. Wanting and expecting share one currency."),
        ("Consequence: ", "achieving a goal and not being surprised are literally the same "
         "quantity. The agent acts to make its own predictions come true."),
    ], size=13, gap=10))
    panel(s, Inches(7.3), Inches(1.5), Inches(5.45), Inches(3.6))
    text(s, Inches(7.65), Inches(1.75), Inches(4.8), Inches(3.2), [
        ("Why this matters for driving", 13, PURPLE, True),
        ("The progress-versus-caution trade-off that traffic psychology has described verbally for "
         "fifty years — zero-risk theory, task-capability interface — needs no arbitration "
         "mechanism here.", 12.5, INK, False, 8),
        ("Goal seeking and caution are one scalar objective, traded in the same units.",
         12.5, INK, False, 8),
    ], spacing=1.15)
    notes(s, """
~3 min. Everyone should get this slide. The one thing to land: preference is encoded as probability, so
'I want' and 'I expect' merge - that is the whole conceptual move. For the HF half, connect to Great
expectations: this is the computational version of the predictive-processing account of driving. No
equations yet; they come next, and only one of them matters.
""")

    # ---- 6 - EFE --------------------------------------------------------------------
    s = head_slide(prs, "The one equation: expected free energy",
                   kicker="Part 2 · Active inference")
    panel(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(2.05))
    text(s, Inches(0.95), Inches(1.72), Inches(11.5), Inches(1.7), [
        ("G(policy)  =  pragmatic value  +  epistemic value", 19, PURPLE, True),
        ("pragmatic:  how well the predicted future matches preferred observations (goal seeking)",
         13, INK, False, 10),
        ("epistemic:  how much the policy is expected to reduce uncertainty (information seeking)",
         13, INK, False, 4),
    ], spacing=1.1)
    text(s, Inches(0.55), Inches(3.95), Inches(12.2), Inches(2.9), bullets([
        ("The agent picks the policy with the lowest expected free energy ", "over a rolling "
         "horizon — caution and progress compete inside one number."),
        ("Epistemic value explains uncertainty-driven behavior for free: ", "slowing at an occluded "
         "junction, looking before pulling out — as information seeking, not a bolted-on rule."),
        ("In collision avoidance the pragmatic term dominates ", "— which is why this talk "
         "mostly needs only the preference function inside it."),
    ], size=14, gap=10))
    notes(s, """
~3 min. Math slide, but keep it verbal. HF half takeaway: two motivations - get what you want, learn what
you need - added in the same units. Analytics half: G(pi) = -E[log p(o)] - expected information gain;
happy to write it out in Q&A. The last bullet is a signpost: for comfort zones we will only need the
pragmatic term, so anyone lost in epistemic value can relax.
""")

    # ---- 7 - machinery ------------------------------------------------------------
    s = head_slide(prs, "Making it computable: the Waymo/TU Delft machinery",
                   kicker="Part 2 · Active inference")
    rows = [
        ("Component", "Choice", "Why"),
        ("Generative model", "discrete-time POMDP, bicycle dynamics", "tractable, interpretable"),
        ("Belief over states", "particle filter, N ≈ 75", "handles multimodality"),
        ("Prediction", "roll particles forward, noise on others' controls",
         "long tail of other-agent behavior"),
        ("Policy search", "CEM model-predictive control", "derivative-free"),
        ("Horizon", "6 s (30 steps at 0.2 s)", "covers an avoidance maneuver"),
        ("EFE estimate", "sample average over particles", "no intractable integrals"),
    ]
    ys = Inches(1.55)
    for r, row in enumerate(rows):
        head = r == 0
        if head:
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), ys + Inches(0.34),
                                     Inches(12.2), Pt(1.6))
            _plain(bar, PURPLE)
        xx = 0.55
        for c, cell in enumerate(row):
            wcol = (2.6, 5.4, 4.2)[c]
            text(s, Inches(xx), ys, Inches(wcol - 0.15), Inches(0.4),
                 [(cell, 12 if not head else 12.5,
                   PURPLE if head else (INK if c < 2 else GREY), head or c == 0)])
            xx += wcol
        ys += Inches(0.42 if not head else 0.48)
    panel(s, Inches(0.55), Inches(4.95), Inches(12.2), Inches(1.7))
    text(s, Inches(0.95), Inches(5.15), Inches(11.5), Inches(1.4), [
        ("Bounded rationality is deliberate", 13.5, PURPLE, True),
        ("The policy budget is capped and the search sometimes returns a sub-optimal plan — on "
         "purpose. Humans are not optimal planners, and a model that plans optimally reproduces the "
         "wrong behavior.", 12.5, INK, False, 5),
    ], spacing=1.15)
    notes(s, """
~2.5 min, aimed at the analytics half. HF half takeaway: the model holds beliefs as a cloud of
hypotheses and plans a few seconds ahead, imperfectly on purpose. The bounded-rationality panel matters
to both audiences: sub-optimality here is a modeling commitment, not a bug. Cost warning for anyone
wanting to run it: roughly 18 s of CPU per simulated timestep.
""")

    # ---- 8 - three additions ---------------------------------------------------------
    s = head_slide(prs, "Three additions make it work for collision avoidance",
                   kicker="Part 2 · Active inference")
    cols = [
        ("Looming perception", BLUE, [
            "The agent observes optical angle and its rate — not distance and speed",
            "Distance-dependent uncertainty and a detection threshold fall out for free",
            ("Consequence: ", "detection delay is derived, not fitted"),
        ]),
        ("Norm-conditioned prediction", TEAL, [
            "Predicted other-agent futures lean toward norm compliance",
            "Trust is capped by currently observed compliance",
            ("Consequence: ", "relaxed in normal traffic, opens to the kinematic long tail the "
             "instant a norm is violated"),
        ]),
        ("Surprise-gated re-planning", PINK, [
            "Plans incrementally by default; full re-plan only when accumulated surprise crosses "
            "a threshold",
            ("Consequence: ", "response timing comes from the model — this is the piece the "
             "comfort-zone method reuses"),
        ]),
    ]
    xs = 0.55
    for heading, colr, items in cols:
        column(s, Inches(xs), Inches(1.55), Inches(3.95), Inches(4.9), heading, items,
               rule=colr, size=12)
        xs += 4.22
    notes(s, """
~3 min. These three are what turn generic active inference into a model of human collision avoidance
(Schumann et al. 2026, Nature Communications). Looming will resonate with the HF half - it is classic
perception work, and the model inherits its consequences instead of fitting them. The third column is
the handoff: say explicitly that the comfort-zone method is built on that surprise signal.
""")

    # ---- 9 - timing mechanism -----------------------------------------------------------
    s = head_slide(prs, "Where response timing comes from",
                   kicker="Part 2 · Active inference")
    panel(s, Inches(0.55), Inches(1.5), Inches(12.2), Inches(1.85))
    text(s, Inches(0.95), Inches(1.7), Inches(11.5), Inches(1.5), [
        ("E(t) = E(t−1) + λ · ε(t)        →  full re-plan when E ≥ 1",
         18, PURPLE, True),
        ("ε(t) is the residual information of the pragmatic value: how far short of the "
         "preferred future the current plan now falls", 13, INK, False, 8),
    ], spacing=1.15)
    text(s, Inches(0.55), Inches(3.7), Inches(12.2), Inches(3.1), bullets([
        ("Zero floor: ", "while the current plan still delivers the preferred future, ε = 0 and "
         "nothing accumulates. An unfolding but still-comfortable situation stays quiet."),
        ("For the human-factors half: ", "this is a drift-diffusion model whose drift rate is "
         "computed from the driver's own generative model, not fitted per scenario."),
        ("For the analytics half: ", "ε depends on the predicted future under the current "
         "policy, so response timing is kinematics-dependent automatically — the empirical "
         "finding (Markkula et al.) that fixed reaction times cannot reproduce."),
    ], size=14, gap=10))
    notes(s, """
~3 min. The most important mechanism slide of the talk - worth slowing down for. The DDM bridge is the
gift to the HF half: accumulate evidence, act at threshold, familiar since the response-time literature -
except the drift rate is supplied by information theory instead of being a free parameter. Land the
zero-floor property hard; the comfort-zone definition in part 4 rests entirely on it.
""")

    # ---- 10 - vs alternatives ---------------------------------------------------------------
    s = head_slide(prs, "What this buys over the alternatives",
                   kicker="Part 2 · Active inference")
    rows = [
        ("", "Mechanistic (TTC rules, looming-DDM)", "Learned behavior models", "Active inference"),
        ("Generalizes across scenarios", "one model per scenario", "yes", "yes — one held out"),
        ("Timing + choice + execution", "usually one of them", "yes", "yes"),
        ("Interpretable", "yes", "no", "beliefs and preferences"),
        ("Safety-critical tail", "if built for it", "under-represented in data",
         "norm-conditioned tail"),
        ("Uncertainty-driven behavior", "mostly not", "implicit", "explicit, epistemic value"),
        ("Cost", "cheap", "data-hungry", "hand-built model, 13 tuned params"),
    ]
    ys = Inches(1.6)
    for r, row in enumerate(rows):
        head = r == 0
        if head:
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), ys + Inches(0.56),
                                     Inches(12.2), Pt(1.6))
            _plain(bar, PURPLE)
        xx = 0.55
        for c, cell in enumerate(row):
            wcol = (3.3, 3.3, 2.8, 2.8)[c]
            text(s, Inches(xx), ys, Inches(wcol - 0.18), Inches(0.55),
                 [(cell, 11.5, PURPLE if head else (INK if c == 0 else GREY), head or c == 0)],
                 spacing=1.0)
            xx += wcol
        ys += Inches(0.68 if head else 0.62)
    text(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.8),
         [("The honest reading: not a better curve fit than a learned model — a unification of "
           "mechanisms that were previously separate, at the price of a hand-built generative model.",
           12.5, PURPLE, True)], spacing=1.1)
    notes(s, """
~2.5 min. Do not oversell. The last row and the closing line are the credibility move: 13 hand-tuned
parameters and a hand-built model is a real cost, and the falsifiability worry is real too - the
strongest counter is the held-out intersection scenario, which the authors did not fit and still
predicted. That held-out logic is exactly the template for our validation study in part 5.
""")

    # ---- PART 3 -----------------------------------------------------------------
    section_slide(prs, 3, "Replication: what worked, what did not, and why it matters where", """
~30 s. Part 3 is the due-diligence part: what was replicated, what was not, and why the failures do not
block the comfort-zone method.
""")

    # ---- 11 - replication ------------------------------------------------------------
    s = head_slide(prs, "Replicating the authors' own code: one match, one diagnosis",
                   kicker="Part 3 · Replication")
    picture(s, "replication_rear_end.png", Inches(7.0), Inches(1.55), w=Inches(5.8))
    text(s, Inches(0.55), Inches(1.5), Inches(6.1), Inches(5.4), bullets([
        ("Setup: ", "rear-end conflict, lead vehicle brakes; the paper's Fig. 3 conditions; their "
         "code, CPU-only, ~18 s per simulated timestep."),
        ("Fig. 3b (25 m/s, 1.0 s gap): ", "reproduced. Response time 0.80 s against a published "
         "0.6–0.8 s, correct brake-and-swerve maneuver."),
        ("Fig. 3a (15 m/s, 1.5 s gap): ", "not reproduced — 0.92 s against 1.4 s, swerves "
         "where the paper only brakes."),
        ("Diagnosed, not mysterious: ", "the calibration table shipped with the code does not span "
         "this operating range, so an assumed-deceleration parameter saturates at its most "
         "pessimistic value. Early-and-evasive at both conditions is exactly what that predicts."),
        ("New — the paper's own runs (OSF deposit, 896 trials): ", "response-time median 1.20 s, "
         "IQR 0.80–1.80 s. Both of our replicated response times sit inside their IQR."),
    ], size=12.5, gap=9))
    notes(s, """
~2.5 min. The message is not 'we matched' but 'we understand the mismatch' - a diagnosed failure with a
mechanism is worth more than an undiagnosed success. Analytics half will appreciate that the failure
mode was predicted by the parameter's direction. The last bullet is new: we now hold the authors' own
simulation output (their OSF deposit, 28 baseline conditions x 32 seeds), so the comparison is against
their true response-time distribution rather than numbers read off figures - and our two points fall
inside their interquartile range. Mention compute honestly: a 10-second scenario is a two-hour CPU job,
which shaped everything we chose to run.
""")

    # ---- 12 - reimplementation ------------------------------------------------------------
    s = head_slide(prs, "An independent re-implementation: right where it rests on preferences",
                   kicker="Part 3 · Replication")
    picture(s, "validation_rear_end.png", Inches(7.0), Inches(1.55), w=Inches(5.8))
    text(s, Inches(0.55), Inches(1.5), Inches(6.1), Inches(5.4), bullets([
        ("Built from the paper and its Supplementary Information ", "in NumPy, to make every "
         "mechanism inspectable; run over the paper's full 28-condition sweep, 140 runs."),
        ("Reproduces 2 of 6 published relations: ", "maneuver choice versus speed, and braking "
         "magnitude versus urgency — precisely the two that rest on the preference function."),
        ("The other four fail on response timing, ", "now quantified against the authors' own runs: "
         "median 1.0 s too early (0.20 s vs their 1.20 s), double the dispersion (sd 1.23 vs "
         "0.66 s), collisions 34% vs their 7%. Diagnosed: the surprise signal is inflated by "
         "control-effort noise in un-smooth planned policies."),
        ("Verdict: ", "the preference function is verified; the closed-loop timing is not. Do not "
         "use it for response-time claims yet."),
    ], size=12.5, gap=9))
    notes(s, """
~2.5 min. The pattern in the successes is the finding: everything resting on the preference function
replicates, everything resting on closed-loop timing does not. That split is what licenses the
comfort-zone method - coming two slides from now - because the method needs only the verified half.
The timing numbers are no longer figure-read: they come from the authors' OSF deposit, so the defect
is measured against the real target distribution.
""")

    # ---- 13 - surprise library -----------------------------------------------------------------
    s = head_slide(prs, "A verified library of surprise measures",
                   kicker="Part 3 · Replication")
    picture(s, "surprise_cutin.png", Inches(7.0), Inches(1.55), w=Inches(5.8))
    text(s, Inches(0.55), Inches(1.5), Inches(6.1), Inches(5.4), bullets([
        ("Three families ", "from Modirshanechi et al.'s taxonomy — probabilistic mismatch, "
         "belief mismatch, observation mismatch — plus the two measures the Waymo group "
         "introduced: residual information and antithesis."),
        ("One interface ", "across Gaussians, mixtures, particle sets and categorical beliefs."),
        ("Tests verify the claims, not just the code: ", "the zero-floor property, invariance to "
         "discretization, the published equivalences between measures, and that antithesis stays "
         "silent on unsurprising belief sharpening."),
        ("Figure: ", "the measures responding to a cut-in — different measures fire at "
         "different moments, which is itself informative."),
    ], size=12.5, gap=9))
    notes(s, """
~2 min, can be compressed if running late. The point for both audiences: 'surprise' is not one thing -
what you compare with what defines a family, and the choice changes when detection fires. Residual
information wins for our purpose on two properties: exactly zero when the expected thing happens, and
no binning parameter on continuous data. Those two properties drive the next part.
""")

    # ---- 14 - separation ---------------------------------------------------------------------
    s = head_slide(prs, "Why none of this blocks the comfort-zone method",
                   kicker="Part 3 · Replication")
    panel(s, Inches(0.55), Inches(1.6), Inches(12.2), Inches(1.95))
    text(s, Inches(0.95), Inches(1.85), Inches(11.4), Inches(1.6), [
        ("The comfort-zone method needs only the preference function and recorded kinematics.",
         17, PURPLE, True),
        ("No closed-loop agent, no particle filter, no policy search, no GPU — so the timing "
         "defect in the re-implementation does not touch it.", 13.5, INK, False, 8),
    ], spacing=1.15)
    text(s, Inches(0.55), Inches(4.0), Inches(12.2), Inches(2.6), bullets([
        ("Verified and used: ", "the preference function (checked term by term against the "
         "Supplementary Information) and residual information (property-tested)."),
        ("Not used: ", "the closed-loop controller whose response-time dispersion is still open."),
        ("This separation is a design decision, ", "made so the method is applicable to "
         "naturalistic data now, at scale, rather than after the controller is fixed."),
    ], size=14, gap=10))
    notes(s, """
~1.5 min. Short but strategically the most important slide in part 3 - it pre-empts the obvious
objection ('your model does not even replicate, why trust the boundary?'). The answer: the boundary
never touches the part that fails. Pause here; make sure this lands before moving to the method itself.
""")

    # ---- PART 4 -----------------------------------------------------------------
    section_slide(prs, 4, "The method: comfort-zone boundaries as level sets", """
~30 s. The core of the talk. Everything so far was ingredients.
""")

    # ---- 15 - the bridge --------------------------------------------------------------
    s = head_slide(prs, "The bridge: every construct has a computational counterpart",
                   kicker="Part 4 · The method")
    rows = [
        ("Traffic-psychology construct", "Active-inference counterpart"),
        ("Comfort zone",
         "region where the current plan still achieves preferred observations (ε ≈ 0)"),
        ("Comfort-zone boundary",
         "level set where ε departs from zero and evidence starts to accumulate"),
        ("Dread-zone boundary",
         "states from which no policy can restore preferred observations"),
        ("Extra motives", "a reshaping of the preference prior p(o)"),
        ("Boundary crossing → evasive action",
         "accumulated evidence crosses threshold → full re-plan"),
    ]
    ys = Inches(1.65)
    for r, row in enumerate(rows):
        head = r == 0
        if head:
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), ys + Inches(0.36),
                                     Inches(12.2), Pt(1.6))
            _plain(bar, PURPLE)
        text(s, Inches(0.55), ys, Inches(4.4), Inches(0.5),
             [(row[0], 13, PURPLE if head else INK, True)])
        text(s, Inches(5.15), ys, Inches(7.6), Inches(0.5),
             [(row[1], 13, PURPLE if head else INK, head)])
        ys += Inches(0.56 if head else 0.72)
    text(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.7),
         [("Not an analogy: the safety-margin term inside the paper's preference function already "
           "is a comfort-zone boundary — the paper just does not call it that.",
           12.5, PURPLE, True)], spacing=1.1)
    notes(s, """
~3 min. Walk the table row by row - this is the slide where the two halves of the audience meet. The
left column is fifty years of traffic psychology; the right column is computable. Close on the bottom
line: the model's own safety-margin preference penalizes exactly the states from which a hard lead-brake
plus a 1-second response would no longer be survivable - a counterfactual escape-route margin, which is
much closer to what the CZB literature means than a bare TTC threshold.
""")

    # ---- 16 - the field --------------------------------------------------------------
    s = head_slide(prs, "The proposal: one field, and the boundary is a level set",
                   kicker="Part 4 · The method")
    picture(s, "czb_field.png", Inches(6.75), Inches(1.5), w=Inches(6.05))
    panel(s, Inches(0.55), Inches(1.5), Inches(5.85), Inches(1.5))
    text(s, Inches(0.9), Inches(1.68), Inches(5.2), Inches(1.2), [
        ("ε(x) = max log p(o) − log p(o(x)) ≥ 0", 15.5, PURPLE, True),
        ("the residual information of the pragmatic value, evaluated at state x",
         12, INK, False, 6),
    ], spacing=1.12)
    text(s, Inches(0.55), Inches(3.3), Inches(5.9), Inches(3.7), bullets([
        ("Comfort zone: ", "ε(x) ≤ c.  Boundary: the level set ε(x) = c."),
        ("Zero floor: ", "ε is exactly 0 at the preferred state — the zone is genuinely "
         "quiet inside."),
        ("Decomposable: ", "log p(o) is a sum of terms (speed, effort, lane, collision, safety "
         "margin), so every exceedance can be attributed to a cause."),
        ("Indicators become projections: ", "min TTC, THW, clearance are slices of this one "
         "surface — which is why they differ per scenario."),
    ], size=12.5, gap=8))
    notes(s, """
~3 min. The definition slide. Figure: the field over the state plane, comfort and dread boundaries drawn
as contours. Three properties to name: zero inside (so exceedance is an event, not a threshold on an
always-positive signal); additive decomposition (so the field says WHY you left the zone); and it is the
same quantity that drives response timing - the boundary and the action mechanism are one object seen
statically and dynamically, not two theories glued together.
""")

    # ---- 17 - closed form + sign error ------------------------------------------------
    s = head_slide(prs, "For car following the boundary has a closed form",
                   kicker="Part 4 · The method")
    text(s, Inches(0.55), Inches(1.5), Inches(6.1), Inches(5.2), bullets([
        ("The safety-margin term asks: ", "if the lead braked hard and I responded after t_react, "
         "would the deceleration I then need exceed a_limit?"),
        ("Setting required = allowed and solving for the gap ", "gives the boundary in closed "
         "form — critical gap or critical THW as a function of the speeds."),
        ("Cross-checked against the numeric level set ", "of the field: agreement to 0.000 m over "
         "45 speeds."),
    ], size=13.5, gap=10))
    panel(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(4.7))
    text(s, Inches(7.35), Inches(1.72), Inches(5.1), Inches(4.3), [
        ("Why the cross-check earns its keep", 13.5, PURPLE, True),
        ("The first comparison disagreed: a sign error on the lead vehicle's own stopping-distance "
         "term had inflated the critical THW from about 0.7 s to about 3.2 s.",
         12.5, INK, False, 8),
        ("A 3.2 s margin is not obviously absurd — the error would plausibly have survived "
         "inspection.", 12.5, INK, False, 8),
        ("Two independent routes to the same surface caught it; a reasonableness check would not "
         "have.", 12.5, PURPLE, True, 8),
    ], spacing=1.18)
    notes(s, """
~2.5 min. The sign-error story is worth telling in full - it is the best 60 seconds of methodology in
the talk. The wrong value was plausible; only computing the same boundary two independent ways exposed
it. For the analytics half this argues for property tests over spot checks; for the HF half it says the
numbers on the next slide have been earned, not eyeballed.
""")

    # ---- 18 - comfort vs dread -----------------------------------------------------------
    s = head_slide(prs, "Comfort versus dread is one interpretable parameter",
                   kicker="Part 4 · The method")
    picture(s, "czb_thw.png", Inches(6.9), Inches(1.55), w=Inches(5.9))
    text(s, Inches(0.55), Inches(1.45), Inches(6.1), Inches(1.95), bullets([
        ("a_limit — the deceleration a driver is prepared to plan around — selects the "
         "boundary:"),
        ("8 m/s²: ", "dread — beyond it, physics allows no escape."),
        ("~4 m/s²: ", "comfort — drivers do not voluntarily plan around emergency "
         "braking."),
    ], size=12.5, gap=6))
    text(s, Inches(0.55), Inches(3.6), Inches(5.9), Inches(0.3),
         [("Critical THW, steady following (lead −6 m/s², t_react 1 s)",
           12, PURPLE, True)])
    table(s, Inches(0.55), Inches(4.0), (1.7, 2.0, 2.1), [
        ("speed", "dread (8 m/s²)", "comfort (4 m/s²)"),
        ("10 m/s", "0.85 s", "1.48 s"),
        ("15 m/s", "0.73 s", "1.67 s"),
        ("20 m/s", "0.61 s", "1.86 s"),
        ("30 m/s", "0.40 s", "2.27 s"),
    ], size=11.5)
    text(s, Inches(0.55), Inches(6.05), Inches(6.1), Inches(1.1),
         [("The comfort boundary lands at 1.5–2.3 s — the range of observed following "
           "headways. A sanity check, not a validation.", 12, GREY, False)], spacing=1.12)
    notes(s, """
~3 min. The two-boundary structure of the CZB literature drops out of one physical parameter, which is
the interpretability argument in miniature. Note the falsifiable prediction hiding in the closed form:
when the assumed lead deceleration equals your own braking limit, critical THW collapses to t_react and
becomes speed-invariant - a strong, clean empirical target. Always quote boundary values together with
the assumed lead deceleration and t_react; they move the boundary and are assumptions, not measurements.
""")

    # ---- 19 - extra motives -------------------------------------------------------------------
    s = head_slide(prs, "Extra motives become measurable parameters",
                   kicker="Part 4 · The method")
    text(s, Inches(0.55), Inches(1.55), Inches(12.2), Inches(1.6), bullets([
        ("In Näätänen and Summala's framing, ", "extra motives make a driver accept "
         "normally unacceptable discomfort — a verbal construct."),
        ("Here they are not a separate mechanism: ", "they are a reshaping of the preference "
         "prior, and the boundary moves as a consequence — by a predicted amount."),
    ], size=14, gap=8))
    table(s, Inches(0.55), Inches(3.5), (2.6, 4.4, 2.2), [
        ("motive", "parameter change", "THW* at 15 m/s"),
        ("baseline", "—", "1.67 s"),
        ("hurried, alert", "t_react 1.0 → 0.6 s", "1.27 s"),
        ("trusting the lead", "assumed lead braking −6 → −3 m/s²", "0.42 s"),
    ], size=12.5, row_h=0.42)
    panel(s, Inches(0.55), Inches(5.55), Inches(12.2), Inches(1.3))
    text(s, Inches(0.95), Inches(5.72), Inches(11.4), Inches(1.0), [
        ("This is the research value: the theory now predicts how much the boundary should move "
         "under a motive manipulation, not just that it moves — and each parameter is "
         "separately measurable.", 13.5, PURPLE, True)], spacing=1.12)
    notes(s, """
~2.5 min. For the HF half this is the payoff slide: a fifty-year-old verbal construct becomes two or
three interpretable, separately measurable parameters. The 'trusting the lead' row is worth dwelling on
- assumed lead capability moves the boundary far more than reaction time does, which itself is a
testable claim. A simulator study manipulating time pressure tests the quantitative prediction directly.
""")

    # ---- 20 - from kinematics --------------------------------------------------------------
    s = head_slide(prs, "The field is computable from recorded kinematics alone",
                   kicker="Part 4 · The method")
    picture(s, "czb_trajectory.png", Inches(6.9), Inches(1.55), w=Inches(5.9))
    text(s, Inches(0.55), Inches(1.5), Inches(6.1), Inches(5.3), bullets([
        ("Evaluating ε along a recorded trajectory ", "needs only the kinematics and the "
         "preference function — no model roll-out, no particle filter, no GPU."),
        ("That is what makes it applicable to naturalistic data at scale, ", "today, despite the "
         "closed-loop model's open issues."),
        ("Figure: ", "a rear-end trajectory; the field's exceedance precedes the driver's brake "
         "onset by ~0.2 s."),
        ("Calibration: ", "the level c is fitted by matching first exceedance to observed response "
         "onset across events — one scalar for all scenarios."),
    ], size=13, gap=10))
    notes(s, """
~2.5 min. The practicality slide. Everything heavy in part 2 - particle filters, policy search - is not
needed here; scoring a trajectory is microseconds, not hours. One honest caveat to voice: the exceedance
in the figure is a model tracing model-consistent behavior, not evidence about humans. It shows the
pipeline runs end to end; the human test is next.
""")

    # ---- 20b - dry run on the authors' data -----------------------------------
    s = head_slide(prs, "First test: it recovers the reference model's response onsets",
                   kicker="Part 4 · The method")
    picture(s, "replication/osf/fig_calibration.png", Inches(6.35), Inches(1.75), w=Inches(6.6))
    text(s, Inches(0.55), Inches(1.5), Inches(5.6), Inches(5.4), bullets([
        ("The paper's OSF deposit ", "holds the authors' own simulation output: 28 rear-end "
         "conditions × 32 seeds = 896 trials, with full kinematics."),
        ("Dry run of the whole pipeline: ", "ε(t) from their kinematics alone, one level c "
         "fitted to their model's brake onsets."),
        ("Result: ", "median onset error 0.0 s (IQR 0.2 s), onset-matching score 0.855 — the "
         "reference model's decisions to act are recoverable from a level set of the field."),
        ("Two caveats: ", "the onsets are the model's, not humans'; and the steep field makes "
         "c weakly identified here — boundary location robust, level not yet disciplined."),
    ], size=12.5, gap=9))
    notes(s, """
~2.5 min. New result. This is the strongest evidence the pipeline works that can exist before human
data: the calibration method, run end to end on the reference implementation's own output, recovers
when that model decides to act - with zero median timing error - from recorded kinematics alone. Be
scrupulous with the caveats: it validates the machinery, not the empirical claim, because the 'driver'
here is the model itself. The weak identification of c is worth saying out loud to the analytics half:
the safety term's step makes the score curve flat in c over orders of magnitude here; scenarios where
the field rises gradually (lateral clearance) will discipline c - another reason the cross-scenario
test matters. Bonus anecdote if time allows: the shortest-gap condition starts 11 cm inside our
closed-form dread boundary, and in the authors' own runs the model refuses it - 72 percent of seeds
leave the road before the lead even brakes. The static boundary and the reference model agree the
state is untenable.
""")

    # ---- PART 5 -----------------------------------------------------------------
    section_slide(prs, 5, "The validation study, and what it needs", """
~30 s. Last part: the experiment that turns the construction into a finding - or falsifies it.
""")

    # ---- 21 - the test -----------------------------------------------------------------
    s = head_slide(prs, "The test: one level c, fitted once, must transfer across scenarios",
                   kicker="Part 5 · Validation")
    steps = [
        ("1", "Take datasets with response onsets in several scenario types — Chalmers LTAP/OD "
              "and pedestrian overtaking are the candidates", BLUE),
        ("2", "Compute ε(t) along every trajectory, from kinematics alone", BLUE),
        ("3", "Fit the boundary level c on one scenario, matching first exceedance to observed "
              "onset", TEAL),
        ("4", "Apply the same c unchanged to the held-out scenario — this step is the "
              "experiment", PURPLE),
        ("5", "Compare against per-scenario indicator thresholds fitted the same way; check the "
              "field's decomposition blames the right term in each scenario", TEAL),
    ]
    ys = Inches(1.55)
    for num, body, colr in steps:
        bar = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), ys, Inches(0.42), Inches(0.42))
        _plain(bar, colr)
        tf = bar.text_frame
        tf.text = num
        tf.paragraphs[0].runs[0].font.size = Pt(15)
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.color.rgb = WHITE if colr != TEAL else INK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        text(s, Inches(1.25), ys + Inches(0.02), Inches(11.4), Inches(0.75),
             [(body, 13.5, INK, num == "4")], spacing=1.08)
        ys += Inches(0.82)
    text(s, Inches(0.55), Inches(5.95), Inches(12.2), Inches(1.0),
         [("If a single c transfers, the level-set formulation is doing real work. If it does not, "
           "the preference function is scenario-specific — which the paper's own SI concedes "
           "is possible, and which would be worth knowing precisely.", 13, PURPLE, True)],
         spacing=1.12)
    notes(s, """
~3 min. Step 4 is the experiment; everything else is preparation. Both outcomes are informative, and
saying so out loud is the falsifiability answer to part 2's criticism - the same held-out logic the
paper itself used with the intersection scenario. Mention the attribution check in step 5: the field
should blame the safety-margin term in following and the lateral term in overtaking, for free.
""")

    # ---- 22 - data --------------------------------------------------------------
    s = head_slide(prs, "What data it takes", kicker="Part 5 · Validation")
    column(s, Inches(0.55), Inches(1.55), Inches(3.9), Inches(4.8),
           "Per event, at ≥ 10 Hz", [
               "Ego position, heading, speed, acceleration, steering",
               "Partner position, heading, speed, acceleration — same frame as ego",
               "Vehicle dimensions, lane geometry",
               ("The label: ", "response onset (brake or steer), extracted the same way as the "
                "published work"),
           ], rule=BLUE, size=12)
    column(s, Inches(4.75), Inches(1.55), Inches(3.9), Inches(4.8), "How much", [
        "30–50 events per scenario type for a stable level fit",
        "At least 2 scenario types, ideally 3 — the generalization claim needs them",
        "10+ events per driver to ask whether c is individual",
    ], rule=TEAL, size=12)
    column(s, Inches(8.95), Inches(1.55), Inches(3.9), Inches(4.8), "Candidate sources", [
        ("Chalmers LTAP/OD: ", "boundaries already quantified conventionally — a direct check"),
        ("Pedestrian overtaking / UDRIVE: ", "different geometry, exactly what the test needs"),
        ("Simulator: ", "the only route that manipulates extra motives rather than inferring them"),
    ], rule=PINK, size=12)
    text(s, Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.6),
         [("The full specification is written to be handed to a data owner: "
           "docs/data_requirements.pdf", 12, GREY, False)], spacing=1.1)
    notes(s, """
~2 min. Keep brisk - the details live in the data-requirements document. Two things to stress: the
onset labels must be extracted the same way as the published work (piecewise-linear fit to speed) so
values are comparable, and the simulator column is not a fallback - it is the only design that tests the
quantitative extra-motives prediction, because motives are manipulated instead of inferred.
""")

    # ---- 23 - limitations -------------------------------------------------------------------
    s = head_slide(prs, "Limitations, stated plainly", kicker="Part 5 · Validation")
    text(s, Inches(0.55), Inches(1.55), Inches(12.2), Inches(5.2), bullets([
        ("No human data yet. ", "Every number shown today is a property of a model. The method is "
         "a construction, not a finding, until the transfer test runs."),
        ("The boundary's location inherits assumptions. ", "Assumed lead deceleration and t_react "
         "move it substantially; absolute values must always be quoted with them."),
        ("The preference function is partly scenario-specific. ", "Lateral preference is defined "
         "per scenario in the paper's SI; a general map-based formulation was not attempted."),
        ("The closed-loop model is not fully validated, ", "so claims coupling boundaries to "
         "predicted response times must wait — the method deliberately does not make them."),
        ("The field has a step at the boundary ", "because the safety term is an indicator — "
         "crisp for defining a boundary, awkward for gradient-based analysis."),
    ], size=13.5, gap=12))
    notes(s, """
~2 min. Volunteer all of it before the Q&A does. The first item is the one to say slowly. If asked why
believe any of this pre-data: the preference function was verified against the SI term by term, the
boundary was computed two independent ways, and the comfort boundary lands in the observed headway range
- necessary conditions, not sufficient ones, and that is exactly why the study in part 5 exists.
""")

    # ---- 24 - take home --------------------------------------------------------------
    s = head_slide(prs, "Take-home")
    msgs = [
        ("One field, not one indicator per scenario",
         "The comfort-zone boundary as a level set of the residual information of pragmatic "
         "value — classic indicators are its projections.", PURPLE),
        ("Old constructs, computational counterparts",
         "Comfort zone, dread zone and extra motives map onto quantities inside a published, "
         "peer-reviewed driver model — extra motives become measurable parameters.", BLUE),
        ("Runs on recorded kinematics alone — and has been run",
         "No model roll-out needed; on the authors' own simulation output it recovers the "
         "reference model's response onsets with zero median timing error.", TEAL),
        ("One falsifiable test decides",
         "Fit the level once, apply it unchanged to a held-out scenario. Either outcome is worth "
         "having.", PINK),
    ]
    ys = Inches(1.5)
    for t1, t2, colr in msgs:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), ys, Inches(0.09), Inches(1.12))
        _plain(bar, colr)
        text(s, Inches(0.95), ys, Inches(11.8), Inches(1.12),
             [(t1, 15, INK, True), (t2, 12.5, GREY, False, 3)], spacing=1.1)
        ys += Inches(1.35)
    notes(s, """
~1.5 min. Four sentences, then stop. If one thing survives the drive home, make it the first: one scalar
field replaces the per-scenario indicator zoo, and the indicators become projections of it.
""")

    # ---- 25 - closing -----------------------------------------------------------------
    s = prs.slides.add_slide(prs.slide_layouts[L_CLOSE])
    if s.shapes.title is not None:
        s.shapes.title.text = "Thank you"
    text(s, Inches(0.55), Inches(5.7), Inches(12.2), Inches(1.2), [
        ("Schumann et al. (2026), Active inference as a model of collision avoidance behavior in "
         "human drivers, Nature Communications 17:5009", 12, WHITE, False),
        ("jonas.bargman@chalmers.se", 12, WHITE, False, 4),
    ])
    clean(s)
    notes(s, """
Q&A prompts to have ready: (1) why not just fit a neural network - interpretability and the
safety-critical tail, slide in part 2; (2) is this falsifiable - the transfer test, either outcome
informative; (3) what about interaction between road users - a known limitation of the model family,
the other agent is non-reactive; (4) computational cost - scoring kinematics is trivial, only the
closed-loop agent is expensive.
""")

    prs.save(str(out))
    print("wrote {} - {} slides".format(out, len(prs.slides._sldIdLst)))


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", type=Path, default=HERE / "czb_talk.pptx")
    args = parser.parse_args()
    build(args.out)


if __name__ == "__main__":
    main()
